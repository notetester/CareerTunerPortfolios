# 스토리보드 PPTX/PDF 빌더.
#   annotate : (구식·비권장) assets/<id>-<t>.png 위에 고정좌표 callout 합성. 좌표식이라 실측 정렬과 어긋남.
#              현재는 tools/capture-screens.mjs 가 DOM 실측 정렬된 화면을 그대로 output/annotated/ 에 캡처한다.
#   pptx     : 16:9 슬라이드 PNG(output/slides/) 를 full-bleed 로 한 장씩 → output/CareerTuner_C_Storyboard.pptx
#   pdf      : 같은 슬라이드 PNG 를 한 페이지씩 → output/CareerTuner_C_Storyboard.pdf (PPTX 와 픽셀 일치)
# usage: python build.py [annotate|pptx|pdf|deck]   (deck = pptx+pdf)
import json, sys, os
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).resolve().parent.parent
SPEC = json.loads((ROOT / "spec" / "c-flow.spec.json").read_text(encoding="utf-8"))
ASSETS = ROOT / "assets"
OUT = ROOT / "output"
ANN = OUT / "annotated"
ANN.mkdir(parents=True, exist_ok=True)

HEX = {"purple":"#6C5CE0","amber":"#B5790F","teal":"#0F8E6E","coral":"#C8501F",
       "blue":"#1F6FC0","pink":"#B83A66","green":"#4E8A18","gray":"#6B6A66"}
CYCLE = ["purple","amber","teal","coral","blue","pink"]
def rgb(c, i=0):
    h = HEX.get(c) or HEX[CYCLE[i % len(CYCLE)]]
    return tuple(int(h[k:k+2],16) for k in (1,3,5))

FONT = "C:/Windows/Fonts/malgunbd.ttf"
def font(sz):
    try: return ImageFont.truetype(FONT, sz)
    except Exception: return ImageFont.load_default()

def annotate():
    for f in SPEC["frames"]:
        for target, fname in (("web", f["web"]), ("app", f["app"])):
            src = ASSETS / fname
            if not src.exists(): continue
            im = Image.open(src).convert("RGBA")
            W, H = im.size
            ov = Image.new("RGBA", (W, H), (0,0,0,0))
            d = ImageDraw.Draw(ov)
            r = max(14, int(W/64)); fb = font(int(r*1.25)); lw = max(3, int(W/360))
            for i, c in enumerate(f.get("callouts", [])):
                if (c.get("target") or "web") != target: continue
                col = rgb(c.get("color"), (c.get("n") or i+1)-1)
                x0,y0 = c["nx"]*W, c["ny"]*H
                x1,y1 = x0 + c["nw"]*W, y0 + c["nh"]*H
                d.rounded_rectangle([x0,y0,x1,y1], radius=max(6,int(W/130)),
                                    outline=col+(255,), width=lw, fill=col+(28,))
                cx, cy = x0, y0
                d.ellipse([cx-r,cy-r,cx+r,cy+r], fill=col+(255,), outline=(255,255,255,255), width=max(2,lw//2))
                d.text((cx,cy), str(c.get("n") or i+1), font=fb, fill=(255,255,255,255), anchor="mm")
            out = Image.alpha_composite(im, ov).convert("RGB")
            out.save(ANN / f"{f['id']}-{target}.png")
            print("annotated", f["id"], target)

def build_pptx():
    # 16:9 슬라이드 이미지(output/slides/slide-NN.png, render-slides.mjs+export.mjs slides 산출)를
    # full-bleed 로 한 장씩 배치한다. 레이아웃은 CSS 가 이미 잡았으므로 PPTX 는 이미지만 깐다.
    from pptx import Presentation
    from pptx.util import Inches
    SW, SH = 13.333, 7.5  # 16:9 == 1280x720
    SLIDES = sorted((OUT / "slides").glob("slide-*.png"))
    if not SLIDES:
        print("!! output/slides/*.png 없음 — 먼저 `node render-slides.mjs && node export.mjs slides` 실행"); return
    prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[6]
    for img in SLIDES:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(img), Inches(0), Inches(0), Inches(SW), Inches(SH))
    out = OUT / "CareerTuner_C_Storyboard.pptx"
    prs.save(str(out)); print("pptx ->", out, f"(slides: {len(prs.slides._sldIdLst)})")

def build_pdf():
    # PPTX 와 동일한 슬라이드 PNG 를 그대로 한 페이지씩 PDF 화 → 두 산출물 화면이 픽셀 단위로 일치.
    SLIDES = sorted((OUT / "slides").glob("slide-*.png"))
    if not SLIDES:
        print("!! output/slides/*.png 없음 — 먼저 `node render-slides.mjs && node export.mjs slides` 실행"); return
    imgs = [Image.open(p).convert("RGB") for p in SLIDES]
    out = OUT / "CareerTuner_C_Storyboard.pdf"
    imgs[0].save(str(out), save_all=True, append_images=imgs[1:], resolution=144.0)
    print("pdf ->", out, f"(pages: {len(imgs)})")

mode = sys.argv[1] if len(sys.argv) > 1 else "deck"
if mode == "annotate": annotate()
if mode in ("pptx","deck","both"): build_pptx()
if mode in ("pdf","deck","both"): build_pdf()
